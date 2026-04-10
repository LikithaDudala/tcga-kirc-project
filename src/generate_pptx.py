"""generate_pptx.py — Generate MS Elevate Capstone Project PowerPoint for TCGA KIRC.

Matches the exact MS Elevate template style:
  - Dark navy header bar with white bold title + Microsoft Elevate logo top-right
  - White content area with black body text
  - Green footer bar
  - Title slide: green left stripe, split left/right with project chart on right
  - Widescreen 16:9  (13.33" x 7.5")

Run:
    python src/generate_pptx.py

Output:
    outputs/TCGA_KIRC_Presentation.pptx
"""

import json
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# ── Paths ─────────────────────────────────────────────────────────────────────
HERE    = Path(__file__).parent.parent
FIGURES = HERE / "outputs" / "figures"
RESULTS = HERE / "outputs" / "results"
OUT     = HERE / "outputs" / "TCGA_KIRC_Presentation.pptx"

# ── Slide dimensions ──────────────────────────────────────────────────────────
W = 13.33   # inches  (16:9 widescreen)
H = 7.5     # inches

# ── Colours  (exact MS Elevate template) ──────────────────────────────────────
NAVY   = RGBColor(0x1E, 0x35, 0x5E)   # dark navy header bar
GREEN  = RGBColor(0x00, 0xA0, 0x50)   # green footer bar / left stripe
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
BLACK  = RGBColor(0x00, 0x00, 0x00)
DARK   = RGBColor(0x1A, 0x1A, 0x2E)   # body text
GRAY   = RGBColor(0x55, 0x55, 0x55)
BLUE   = RGBColor(0x00, 0x53, 0xA5)   # section header accent on slides
TEAL   = RGBColor(0x00, 0x70, 0x80)   # title slide right panel
LBKG   = RGBColor(0xEB, 0xF1, 0xFA)   # light blue table background

# Microsoft Windows logo square colours
_MS_R = RGBColor(0xF2, 0x50, 0x22)
_MS_G = RGBColor(0x7F, 0xBA, 0x00)
_MS_B = RGBColor(0x00, 0xA4, 0xEF)
_MS_Y = RGBColor(0xFF, 0xB9, 0x00)


# ── Low-level helpers ─────────────────────────────────────────────────────────

def blank_slide(prs):
    layout = next((l for l in prs.slide_layouts if l.name == "Blank"), prs.slide_layouts[-1])
    return prs.slides.add_slide(layout)


def rect(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def txbox(slide, left, top, width, height):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    box.fill.background()
    tf = box.text_frame
    tf.word_wrap = True
    return tf


def ms_elevate_logo(slide, left, top, on_dark=True):
    """Draw Windows 4-square logo + 'Microsoft / Elevate' text."""
    sq, gap = 0.14, 0.025
    rect(slide, left,            top,            sq, sq, _MS_R)
    rect(slide, left + sq + gap, top,            sq, sq, _MS_G)
    rect(slide, left,            top + sq + gap, sq, sq, _MS_B)
    rect(slide, left + sq + gap, top + sq + gap, sq, sq, _MS_Y)

    txt_color = WHITE if on_dark else DARK
    tx = 2 * sq + gap + 0.07
    tf1 = txbox(slide, left + tx, top - 0.01, 1.5, 0.22)
    r1 = tf1.paragraphs[0].add_run()
    r1.text = "Microsoft"
    r1.font.size = Pt(10)
    r1.font.color.rgb = txt_color

    tf2 = txbox(slide, left + tx, top + 0.18, 1.5, 0.22)
    r2 = tf2.paragraphs[0].add_run()
    r2.text = "Elevate"
    r2.font.size = Pt(10)
    r2.font.color.rgb = txt_color


def add_header(slide, title_text):
    """Navy header bar + white bold title + MS Elevate logo (top-right, on dark)."""
    rect(slide, 0, 0, W, 1.3, NAVY)
    tf = txbox(slide, 0.35, 0.18, W - 2.6, 0.95)
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title_text
    run.font.size = Pt(30)
    run.font.bold = True
    run.font.color.rgb = WHITE
    ms_elevate_logo(slide, W - 2.1, 0.22, on_dark=True)


def add_footer(slide):
    """Green footer bar."""
    rect(slide, 0, H - 0.33, W, 0.33, GREEN)


def content_text(slide, items, left=0.35, top=1.42, width=None, height=None):
    """Render structured bullet list into a slide text box."""
    if width is None:
        width = W - 0.7
    if height is None:
        height = H - top - 0.45
    tf = txbox(slide, left, top, width, height)
    first = True
    for item in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        if isinstance(item, dict) and "h" in item:
            run = p.add_run()
            run.text = item["h"]
            run.font.size = Pt(15)
            run.font.bold = True
            run.font.color.rgb = NAVY
            p.space_before = Pt(9)
        elif isinstance(item, tuple):
            level, text = item
            p.level = min(level, 4)
            run = p.add_run()
            run.text = text
            run.font.size = Pt(13)
            run.font.color.rgb = BLACK
        else:
            run = p.add_run()
            run.text = str(item)
            run.font.size = Pt(13)
            run.font.color.rgb = BLACK


# ── Slide builders ────────────────────────────────────────────────────────────

def slide_01_title(prs):
    """Title slide — matches MS Elevate template: green left stripe, student info, chart panel."""
    slide = blank_slide(prs)

    # White left background
    rect(slide, 0, 0, 8.6, H, WHITE)
    # Green left stripe
    rect(slide, 0, 0, 0.22, H, GREEN)
    # Teal right panel (chart area)
    rect(slide, 8.6, 0, W - 8.6, H, TEAL)
    # Green bottom bar (full width)
    rect(slide, 0, H - 0.33, W, 0.33, GREEN)

    # MS Elevate logo — top-left on white
    ms_elevate_logo(slide, 0.38, 0.25, on_dark=False)

    # "CAPSTONE PROJECT" label
    tf_cap = txbox(slide, 0.45, 1.45, 7.8, 0.42)
    r_cap = tf_cap.paragraphs[0].add_run()
    r_cap.text = "CAPSTONE PROJECT"
    r_cap.font.size = Pt(14)
    r_cap.font.bold = True
    r_cap.font.color.rgb = DARK

    # Project title
    tf_tit = txbox(slide, 0.45, 1.88, 7.8, 1.7)
    tf_tit.word_wrap = True
    p1 = tf_tit.paragraphs[0]
    r1 = p1.add_run()
    r1.text = "TCGA KIRC:"
    r1.font.size = Pt(46)
    r1.font.bold = True
    r1.font.color.rgb = BLACK
    p2 = tf_tit.add_paragraph()
    r2 = p2.add_run()
    r2.text = "Survival Analysis"
    r2.font.size = Pt(38)
    r2.font.bold = True
    r2.font.color.rgb = BLACK

    # "PRESENTED BY" section
    tf_by = txbox(slide, 0.45, 3.75, 8.0, 3.2)
    tf_by.word_wrap = True
    details = [
        ("PRESENTED BY",  None),
        ("STUDENT NAME:", "Likitha Dudala"),
        ("COLLEGE NAME:", "Sardar Vallabhbhai Patel Institute of Technology"),
        ("DEPARTMENT:",   "Computer Engineering"),
        ("EMAIL ID:",     "likithadudala04@gmail.com"),
    ]
    first = True
    for label, value in details:
        p = tf_by.paragraphs[0] if first else tf_by.add_paragraph()
        first = False
        r_lbl = p.add_run()
        r_lbl.text = label + ("" if not value else "  ")
        r_lbl.font.size = Pt(13) if value else Pt(14)
        r_lbl.font.bold = True
        r_lbl.font.color.rgb = DARK
        if value:
            r_val = p.add_run()
            r_val.text = value
            r_val.font.size = Pt(13)
            r_val.font.bold = False
            r_val.font.color.rgb = GRAY
        p.space_before = Pt(4)

    # Right panel image — KM overall survival curve
    img = FIGURES / "03_km_overall.png"
    if img.exists():
        slide.shapes.add_picture(str(img), Inches(8.65), Inches(0.35), Inches(4.5), Inches(6.8))


def slide_02_outline(prs):
    slide = blank_slide(prs)
    add_header(slide, "OUTLINE:")
    add_footer(slide)
    content_text(slide, [
        (0, "Problem Statement"),
        (0, "Proposed System / Solution"),
        (0, "System Development Approach  (Technology Used)"),
        (0, "Algorithm & Deployment"),
        (0, "Result  (Output Images)"),
        (0, "Conclusion"),
        (0, "Future Scope"),
        (0, "References"),
    ], top=1.55)


def slide_03_problem(prs):
    slide = blank_slide(prs)
    add_header(slide, "PROBLEM STATEMENT:")
    add_footer(slide)
    content_text(slide, [
        (0, "Kidney Renal Clear Cell Carcinoma (KIRC) is the most prevalent kidney cancer subtype, "
            "accounting for ~75% of all kidney malignancies. It has a 5-year survival rate below "
            "12% for advanced metastatic disease, making accurate prognosis vital."),
        (0, ""),
        (0, "Despite sharing the same histological type, KIRC patients exhibit dramatically "
            "different survival outcomes. Current AJCC clinical staging (Stages I–IV) is "
            "insufficient to predict individual patient prognosis or guide personalised treatment."),
        (0, ""),
        (0, "The Cancer Genome Atlas (TCGA) provides high-dimensional genomic and clinical data for "
            "533 KIRC patients. However, right-censored survival data and the high-dimensional gene "
            "expression space (20,530 genes) make standard regression methods inappropriate."),
        (0, ""),
        (0, "There is a critical need for robust, interpretable, data-driven survival models that "
            "can: (1) accurately stratify patients into high- and low-risk groups, and "
            "(2) identify the key genomic and clinical features driving mortality in KIRC."),
    ])


def slide_04_solution(prs):
    slide = blank_slide(prs)
    add_header(slide, "PROPOSED SOLUTION:")
    add_footer(slide)
    content_text(slide, [
        (0, "The proposed system builds an end-to-end survival analysis pipeline on TCGA-KIRC data, "
            "applying four complementary machine learning models to predict patient survival and "
            "identify key prognostic biomarkers."),
        (0, ""),
        {"h": "Data Collection"},
        (1, "TCGA-KIRC clinical.tsv + follow_up.tsv  (GDC Portal, 537 patients, OS labels)"),
        (1, "kirc_expression.tsv via UCSC Xena  (20,530 genes x 606 tumour samples)"),

        {"h": "Data Preprocessing"},
        (1, "Survival time and event derivation  |  patient ID alignment across datasets"),
        (1, "Variance-based gene filtering: top 2,000 genes retained  |  StandardScaler normalisation"),

        {"h": "Machine Learning Algorithms"},
        (1, "Cox PH  |  LASSO-penalised Cox  |  Random Survival Forest  |  DeepSurv (PyTorch MLP)"),
        (1, "70/30 stratified train/test split  |  5-fold CV for hyperparameter tuning"),

        {"h": "Deployment"},
        (1, "Live Streamlit dashboard: https://tcga-kirc.streamlit.app/  |  "
            "GitHub: https://github.com/LikithaDudala/tcga-kirc-project"),

        {"h": "Evaluation"},
        (1, "Concordance Index (C-index), Integrated Brier Score, Kaplan-Meier risk stratification"),
    ])


def slide_05_system(prs):
    slide = blank_slide(prs)
    add_header(slide, "SYSTEM DEVELOPMENT APPROACH  (Technology Used):")
    add_footer(slide)

    # Left column
    content_text(slide, [
        {"h": "System Requirements"},
        (1, "Python 3.12  |  Jupyter Notebook"),
        (1, "Git + GitHub  (version control & hosting)"),
        (1, "Streamlit Cloud  (live deployment)"),
        (0, ""),
        {"h": "Data Sources"},
        (1, "TCGA GDC Portal  —  clinical TSVs"),
        (1, "UCSC Xena  —  RNA-seq gene expression"),
        (1, "529 patients  |  20,530 genes  ->  2,000"),
    ], left=0.35, top=1.42, width=6.1, height=5.7)

    # Right column
    content_text(slide, [
        {"h": "Libraries Required to Build the Model"},
        (1, "pandas, numpy          —  data wrangling"),
        (1, "lifelines              —  Cox PH, Kaplan-Meier"),
        (1, "scikit-survival        —  LASSO Cox, RSF"),
        (1, "PyTorch                —  DeepSurv neural network"),
        (1, "plotly, matplotlib     —  visualisations"),
        (1, "scikit-learn           —  preprocessing & metrics"),
        (1, "streamlit              —  interactive dashboard"),
        (1, "python-pptx            —  this presentation"),
    ], left=6.6, top=1.42, width=6.4, height=5.7)

    # Thin vertical divider
    div = slide.shapes.add_shape(1, Inches(6.45), Inches(1.48), Inches(0.04), Inches(5.6))
    div.fill.solid()
    div.fill.fore_color.rgb = RGBColor(0xCC, 0xD6, 0xE8)
    div.line.fill.background()


def slide_06_algorithm(prs):
    slide = blank_slide(prs)
    add_header(slide, "ALGORITHM & DEPLOYMENT:")
    add_footer(slide)
    content_text(slide, [
        {"h": "Algorithm Selection"},
        (1, "Four complementary survival models were chosen to capture linear, penalised, "
            "ensemble, and deep-learning representations of the survival task:"),
        (1, "Cox Proportional Hazards — clinical features only; interpretable hazard ratios (baseline)"),
        (1, "LASSO-penalised Cox (BEST) — clinical + 2,000 genes; L1 shrinkage auto-selects genes"),
        (1, "Random Survival Forest  — 200 trees; captures non-linear gene interactions"),
        (1, "DeepSurv (PyTorch)  — 3-layer MLP; Cox partial-likelihood loss; Adam; 200 epochs"),

        {"h": "Data Input"},
        (1, "Features: age, AJCC stage, gender + top 2,000 variance-ranked genes"),
        (1, "Target: right-censored overall survival (time in days, event = 1 if deceased)"),

        {"h": "Training Process"},
        (1, "70/30 stratified train/test split  (370 train, 159 test)"),
        (1, "5-fold CV for LASSO alpha  |  batch-norm + dropout(0.3) for DeepSurv"),

        {"h": "Deployment"},
        (1, "Models serialised via joblib (Cox/RSF) and torch.save (DeepSurv)"),
        (1, "Streamlit app loads Cox PH + scaler; user enters clinical inputs -> survival curve"),
        (1, "Live:  https://tcga-kirc.streamlit.app/"),
    ])


def slide_07_results(prs, model_results, cohort):
    slide = blank_slide(prs)
    add_header(slide, "RESULT  (Output Images):")
    add_footer(slide)

    # Model performance table
    rect(slide, 0.3, 1.42, 3.9, 2.35, LBKG)
    tf_tbl = txbox(slide, 0.4, 1.48, 3.7, 2.25)
    hdr = tf_tbl.paragraphs[0]
    r_hdr = hdr.add_run()
    r_hdr.text = "C-index  (Test Set)"
    r_hdr.font.size = Pt(12)
    r_hdr.font.bold = True
    r_hdr.font.color.rgb = NAVY
    rows = [
        ("LASSO Cox             ", model_results["LASSO Cox"],              True),
        ("Cox PH (clinical)     ", model_results["Cox PH (clinical)"],      False),
        ("DeepSurv              ", model_results["DeepSurv"],               False),
        ("Random Survival Forest", model_results["Random Survival Forest"], False),
    ]
    for name, score, best in rows:
        p = tf_tbl.add_paragraph()
        r = p.add_run()
        badge = " *" if best else "  "
        r.text = f" {badge} {name}  {score:.4f}"
        r.font.size = Pt(11)
        r.font.bold = best
        r.font.color.rgb = BLUE if best else DARK

    # Cohort summary
    tf_stats = txbox(slide, 0.4, 3.9, 3.7, 2.8)
    sh = tf_stats.paragraphs[0].add_run()
    sh.text = "Cohort Summary"
    sh.font.size = Pt(12)
    sh.font.bold = True
    sh.font.color.rgb = NAVY
    stats = [
        f"Patients:          {cohort['total_patients']}",
        f"Events (deceased): {cohort['events_dead']}  ({cohort['events_dead'] / cohort['total_patients'] * 100:.1f}%)",
        f"Censored (alive):  {cohort['censored_alive']}",
        f"Median follow-up:  {cohort['median_time_days']:.0f} days",
        f"Genes analysed:    {cohort['num_genes']:,}",
        f"Brier Score (RSF): {cohort['ibs_rsf']:.3f}",
    ]
    for s in stats:
        ps = tf_stats.add_paragraph()
        rs = ps.add_run()
        rs.text = f"  {s}"
        rs.font.size = Pt(11)
        rs.font.color.rgb = DARK

    # 3 output figure screenshots
    figs = [
        ("09_model_comparison.png",   4.35, 1.42, 4.45, 2.9),
        ("03_km_overall.png",         9.0,  1.42, 4.15, 2.9),
        ("10_risk_stratification.png", 4.35, 4.45, 8.8,  2.65),
    ]
    for fname, l, t, w, h_img in figs:
        fpath = FIGURES / fname
        if fpath.exists():
            slide.shapes.add_picture(str(fpath), Inches(l), Inches(t), Inches(w), Inches(h_img))


def slide_08_conclusion(prs, model_results, cohort):
    slide = blank_slide(prs)
    add_header(slide, "CONCLUSION:")
    add_footer(slide)
    delta = model_results["LASSO Cox"] - model_results["Cox PH (clinical)"]
    content_text(slide, [
        (0, "This project developed and evaluated a complete end-to-end survival analysis pipeline "
            "on 529 TCGA-KIRC patients using high-dimensional gene expression data combined with "
            "clinical variables, comparing four distinct model families."),
        (0, ""),
        {"h": "Key Findings"},
        (1, f"LASSO-penalised Cox achieved the best C-index of {model_results['LASSO Cox']:.4f}, "
            f"outperforming clinical-only Cox PH by delta-C = {delta:.4f}."),
        (1, f"Random Survival Forest confirmed non-linear gene interactions "
            f"(Integrated Brier Score = {cohort['ibs_rsf']:.3f})."),
        (1, "Kaplan-Meier risk stratification shows clear high/low-risk separation "
            "(log-rank p < 0.001)."),
        (1, "LASSO identified a sparse, interpretable gene signature predictive of KIRC mortality."),
        (0, ""),
        {"h": "Project Outcomes"},
        (1, "Reproducible pipeline: raw TCGA data -> trained models -> interactive dashboard."),
        (1, "Live Streamlit app: https://tcga-kirc.streamlit.app/"),
        (1, "Source code: https://github.com/LikithaDudala/tcga-kirc-project"),
        (0, ""),
        {"h": "Limitations"},
        (1, "TCGA data may carry selection bias; model not validated on an external cohort."),
    ])


def slide_09_future(prs):
    slide = blank_slide(prs)
    add_header(slide, "FUTURE SCOPE:")
    add_footer(slide)
    content_text(slide, [
        {"h": "Model Improvements"},
        (1, "Ensemble / stacking of all 4 models for improved concordance."),
        (1, "Attention-based DeepSurv with pathway-level gene embeddings (Gene Ontology / KEGG)."),
        (1, "Multi-omics integration: copy number variation, methylation, proteomics."),
        (1, "SHAP values for per-patient risk factor explanation (explainable AI)."),
        (0, ""),
        {"h": "Validation & Clinical Translation"},
        (1, "External validation on ICGC or CheckMate-025 / IMmotion151 datasets."),
        (1, "Prospective clinical pilot: integrate risk scores with EHR systems."),
        (0, ""),
        {"h": "Engineering & Deployment  (App already live: https://tcga-kirc.streamlit.app/)"},
        (1, "REST API (FastAPI) for programmatic risk score queries from clinical systems."),
        (1, "Automated retraining pipeline as new TCGA data releases become available."),
        (1, "Multi-cancer extension: apply pipeline to TCGA-LUAD, TCGA-BRCA, TCGA-GBM."),
    ])


def slide_10_references(prs):
    slide = blank_slide(prs)
    add_header(slide, "REFERENCES:")
    add_footer(slide)
    content_text(slide, [
        {"h": "Datasets"},
        (1, "1.  The Cancer Genome Atlas (TCGA) — GDC Data Portal"),
        (2, "https://portal.gdc.cancer.gov/projects/TCGA-KIRC"),
        (1, "2.  UCSC Xena — RNA-seq gene expression"),
        (2, "https://xenabrowser.net/datapages/"),
        (0, ""),
        {"h": "Libraries & Frameworks"},
        (1, "3.  Davidson-Pilon C. et al. (2019). lifelines. Journal of Open Source Software."),
        (1, "4.  Poelsterl S. (2020). scikit-survival. JMLR 21(212):1-6."),
        (1, "5.  Katzman J.L. et al. (2018). DeepSurv. BMC Med Res Methodol 18:24."),
        (1, "6.  Ricketts C.J. et al. (2018). TCGA KIRC. Cell Reports 23(1):313-326."),
        (0, ""),
        {"h": "Project Links"},
        (1, "GitHub:   https://github.com/LikithaDudala/tcga-kirc-project"),
        (1, "Live App: https://tcga-kirc.streamlit.app/"),
    ])


def slide_11_thankyou(prs):
    """Thank You slide — navy header, white centre, green footer (matches template)."""
    slide = blank_slide(prs)

    # Navy header bar
    rect(slide, 0, 0, W, 1.1, NAVY)
    # MS Elevate logo (top-right, on dark)
    ms_elevate_logo(slide, W - 2.1, 0.2, on_dark=True)
    # Green footer
    rect(slide, 0, H - 0.33, W, 0.33, GREEN)

    # "Thank You" — large, dark navy, centred on white
    tf = txbox(slide, 0.5, 2.0, W - 1.0, 2.2)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "Thank You"
    run.font.size = Pt(64)
    run.font.bold = True
    run.font.color.rgb = NAVY


# ── Main ──────────────────────────────────────────────────────────────────────

def main():
    with open(RESULTS / "model_results.json") as f:
        model_results = json.load(f)
    with open(RESULTS / "cohort_summary.json") as f:
        cohort = json.load(f)

    prs = Presentation()
    prs.slide_width  = Inches(W)
    prs.slide_height = Inches(H)

    slide_01_title(prs)
    slide_02_outline(prs)
    slide_03_problem(prs)
    slide_04_solution(prs)
    slide_05_system(prs)
    slide_06_algorithm(prs)
    slide_07_results(prs, model_results, cohort)
    slide_08_conclusion(prs, model_results, cohort)
    slide_09_future(prs)
    slide_10_references(prs)
    slide_11_thankyou(prs)

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT))
    print(f"Saved -> {OUT}")
    print(f"Slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
